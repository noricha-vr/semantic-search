"""Office文書プロセッサ。

Word、Excel、PowerPointからテキストを抽出する。
"""

from dataclasses import dataclass
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

from src.config.logging import get_logger

logger = get_logger()


@dataclass
class OfficeResult:
    """Office文書処理結果。"""

    text: str
    doc_type: str
    sheet_count: int | None = None  # Excel用
    slide_count: int | None = None  # PowerPoint用
    paragraph_count: int | None = None  # Word用


class OfficeProcessor:
    """Office文書プロセッサ。"""

    DOCX_EXTENSIONS = {".docx", ".doc"}
    XLSX_EXTENSIONS = {".xlsx", ".xls"}
    PPTX_EXTENSIONS = {".pptx", ".ppt"}

    def extract_from_docx(self, file_path: Path | str) -> OfficeResult:
        """Wordドキュメントからテキストを抽出。

        Args:
            file_path: ファイルパス

        Returns:
            テキストと情報
        """
        file_path = Path(file_path)
        doc = DocxDocument(str(file_path))

        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # テーブルからもテキストを抽出
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_parts.append(" | ".join(row_text))

        full_text = "\n\n".join(text_parts)

        logger.info(
            f"Extracted text from Word: {file_path}, "
            f"paragraphs: {len(doc.paragraphs)}"
        )

        return OfficeResult(
            text=full_text,
            doc_type="docx",
            paragraph_count=len(doc.paragraphs),
        )

    def extract_from_xlsx(self, file_path: Path | str) -> OfficeResult:
        """Excelワークブックからテキストを抽出。

        Args:
            file_path: ファイルパス

        Returns:
            テキストと情報
        """
        file_path = Path(file_path)
        wb = load_workbook(str(file_path), data_only=True)

        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]

            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                if row_values:
                    sheet_text.append(" | ".join(row_values))

            if len(sheet_text) > 1:  # シート名以外にデータがある
                text_parts.append("\n".join(sheet_text))

        full_text = "\n\n".join(text_parts)

        logger.info(
            f"Extracted text from Excel: {file_path}, sheets: {len(wb.sheetnames)}"
        )

        return OfficeResult(
            text=full_text,
            doc_type="xlsx",
            sheet_count=len(wb.sheetnames),
        )

    def extract_from_pptx(self, file_path: Path | str) -> OfficeResult:
        """PowerPointプレゼンテーションからテキストを抽出。

        Args:
            file_path: ファイルパス

        Returns:
            テキストと情報
        """
        file_path = Path(file_path)
        prs = Presentation(str(file_path))

        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if len(slide_text) > 1:  # スライド番号以外にテキストがある
                text_parts.append("\n".join(slide_text))

        full_text = "\n\n".join(text_parts)

        logger.info(
            f"Extracted text from PowerPoint: {file_path}, "
            f"slides: {len(prs.slides)}"
        )

        return OfficeResult(
            text=full_text,
            doc_type="pptx",
            slide_count=len(prs.slides),
        )

    def extract_text(self, file_path: Path | str) -> OfficeResult:
        """Office文書からテキストを抽出。

        Args:
            file_path: ファイルパス

        Returns:
            テキストと情報

        Raises:
            FileNotFoundError: ファイルが見つからない
            ValueError: サポートされていない形式
        """
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = file_path.suffix.lower()

        if suffix in self.DOCX_EXTENSIONS:
            return self.extract_from_docx(file_path)
        elif suffix in self.XLSX_EXTENSIONS:
            return self.extract_from_xlsx(file_path)
        elif suffix in self.PPTX_EXTENSIONS:
            return self.extract_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported format: {suffix}")

    def is_supported(self, file_path: Path | str) -> bool:
        """ファイルがサポートされているかを判定。

        Args:
            file_path: ファイルパス

        Returns:
            サポートされていればTrue
        """
        file_path = Path(file_path)
        suffix = file_path.suffix.lower()
        return suffix in (
            self.DOCX_EXTENSIONS | self.XLSX_EXTENSIONS | self.PPTX_EXTENSIONS
        )
